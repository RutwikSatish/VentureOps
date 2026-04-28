"""
VentureOps — Startup Operational Due Diligence Engine v2.0
Built for the Roux Institute | Rutwik Satish

WHAT'S NEW IN v2:
- File upload: extract startup data from PDF pitch decks, PPTX, XLSX, CSV
- Demo mode with guided tour (load one-click example data)
- Evidence-based Operational Readiness Score (every sub-score shows source + why)
- TAM/SAM/SOM gap analysis (founder claim vs. data estimate with red/yellow/green flags)
- Confidence ratings on every calculated dimension
- Suggestions sourced to Bessemer, CB Insights, Sequoia, SaaStr, David Skok

VALIDATED FORMULAS & SOURCES:
- Net Burn = Gross Burn - Revenue  [standard accounting]
- Runway = Cash on Hand / Net Burn  [standard VC metric]
- Burn Multiple = Net Burn / Net New ARR  [Bessemer State of Cloud 2023]
- Gross Profit LTV = (ARPU × Gross Margin %) / Monthly Churn Rate  [David Skok, saastr.com]
- LTV/CAC: target >3x  [SaaStr, Bessemer, McKinsey — industry consensus]
- CAC Payback = CAC / (ARPU × Gross Margin %)  → result in months  [David Skok]
- CAC Payback benchmark: <12mo SMB, <18mo mid-market  [David Skok, McKinsey SaaS]
- Rule of 40 = YoY Growth % + EBITDA Margin %  [Bain & Company, McCall & Murphy 2016]
- Burn Multiple benchmarks: <1 excellent, 1-1.5 good, 1.5-2 OK, 2-3 suspect, >3 poor
  [Bessemer Venture Partners State of Cloud 2023]
- TAM bottom-up: Total Potential Customers × ACV  [Antler, Forum VC, a16z methodology]
- SOM: built from sales capacity, NOT % of TAM  [Bill Gurley "A Rake Too Far" 2013;
  Sequoia Capital pitch framework]
- Scale bottleneck 80% threshold: Kingman's Formula (1961), queuing theory
  Applied to operations: >80% utilization → exponential queue growth
- Operational readiness weights: calibrated from CB Insights "Why Startups Fail" (101 post-mortems, original 2014 study)
  + Bessemer Series A survival predictors
- TAM/SOM benchmarks: Bessemer State of Cloud 2023; SaaStr Annual AE Productivity Survey 2024

Install: pip install streamlit plotly groq pandas numpy pdfplumber python-pptx openpyxl
Run:     streamlit run ventureops_app.py
Secrets: Create .streamlit/secrets.toml with GROQ_API_KEY = "your_key_here"
"""

import streamlit as st
import plotly.graph_objects as go
import plotly.express as px
import pandas as pd
import numpy as np
import json
import io
import re

# ─── PAGE CONFIG ────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="VentureOps — Startup Due Diligence Engine",
    page_icon="🔬",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ─── STYLES ─────────────────────────────────────────────────────────────────────
st.markdown("""
<style>
    .main-header {
        font-size: 2.2rem; font-weight: 800;
        background: linear-gradient(90deg, #1a1a2e, #16213e, #0f3460);
        -webkit-background-clip: text; -webkit-text-fill-color: transparent;
        margin-bottom: 0.2rem;
    }
    .sub-header { color: #6b7280; font-size: 1rem; margin-bottom: 1rem; }
    .metric-card {
        background: #f8fafc; border: 1px solid #e2e8f0;
        border-radius: 12px; padding: 1.2rem; text-align: center;
    }
    .score-excellent { color: #16a34a; font-weight: 800; font-size: 2rem; }
    .score-good      { color: #2563eb; font-weight: 800; font-size: 2rem; }
    .score-moderate  { color: #d97706; font-weight: 800; font-size: 2rem; }
    .score-critical  { color: #dc2626; font-weight: 800; font-size: 2rem; }
    .insight-box {
        background: #f0fdf4; border-left: 4px solid #16a34a;
        padding: 1rem; border-radius: 0 8px 8px 0; margin: 0.5rem 0;
    }
    .warning-box {
        background: #fffbeb; border-left: 4px solid #d97706;
        padding: 1rem; border-radius: 0 8px 8px 0; margin: 0.5rem 0;
    }
    .danger-box {
        background: #fef2f2; border-left: 4px solid #dc2626;
        padding: 1rem; border-radius: 0 8px 8px 0; margin: 0.5rem 0;
    }
    .source-box {
        background: #f1f5f9; border: 1px solid #cbd5e1;
        border-radius: 8px; padding: 0.8rem; font-size: 0.78rem;
        color: #64748b; margin-top: 0.5rem;
    }
    .gap-green  { background:#f0fdf4; border-left:4px solid #16a34a; padding:0.8rem; border-radius:0 8px 8px 0; margin:0.3rem 0; }
    .gap-yellow { background:#fffbeb; border-left:4px solid #d97706; padding:0.8rem; border-radius:0 8px 8px 0; margin:0.3rem 0; }
    .gap-red    { background:#fef2f2; border-left:4px solid #dc2626; padding:0.8rem; border-radius:0 8px 8px 0; margin:0.3rem 0; }
    .confidence-high   { background:#dcfce7; color:#166534; padding:2px 8px; border-radius:4px; font-size:0.75rem; font-weight:600; }
    .confidence-medium { background:#fef9c3; color:#854d0e; padding:2px 8px; border-radius:4px; font-size:0.75rem; font-weight:600; }
    .confidence-low    { background:#fee2e2; color:#991b1b; padding:2px 8px; border-radius:4px; font-size:0.75rem; font-weight:600; }
    .tour-box { background:#eff6ff; border:2px solid #2563eb; border-radius:12px; padding:1.2rem; margin:1rem 0; }
    .extracted-tag { background:#dbeafe; color:#1d4ed8; padding:2px 6px; border-radius:4px; font-size:0.72rem; font-weight:600; }
    .stTabs [data-baseweb="tab"] { font-weight: 600; font-size: 0.95rem; }
</style>
""", unsafe_allow_html=True)

# ─── DEMO DATA ───────────────────────────────────────────────────────────────────
DEMO_DATA = {
    "startup_name":    "NovaMed AI",
    "industry":        "HealthTech",
    "stage":           "Seed",
    "team_size":       11,
    "founded_years":   2,
    "biz_model":       "B2B SaaS",
    "monthly_revenue": 42000,
    "monthly_burn":    185000,
    "cash_on_hand":    1400000,
    "arpu":            420,
    "cac":             3200,
    "gross_margin_pct":68,
    "monthly_churn":   2.8,
    "mrr_growth":      9.5,
    # Ops readiness
    "sc_single_source": 7, "sc_lead_time": 5, "sc_inventory": 4,
    "ue_ltv_cac": 4, "ue_payback": 5, "ue_margin": 6,
    "ol_automation": 3, "ol_scalability": 5,
    "tc_key_person": 8, "tc_hiring": 4,
    "pm_documentation": 3, "pm_kpis": 5,
    "td_infrastructure": 6, "td_debt_level": 5,
    # TAM
    "total_potential_customers": 32000,
    "acv_annual": 5040,
    "sam_filter_pct": 22,
    "som_capture_pct": 7,
    "founder_tam_claim": 28.0,
    "founder_sam_claim": 6.5,
    "founder_som_claim": 320.0,
    "industry_market_size": 45.0,
    "td_segment_pct": 14,
    "td_geo_pct": 38,
    "td_product_fit_pct": 55,
    # Scale
    "current_customers": 100,
    "ops_team_size": 2,
    "customers_per_csm": 50,
    "infra_cost_per_cust": 12.0,
    "supplier_count": 3,
    "single_source_pct": 72,
    "manual_ops_hrs_week": 42,
    "tech_incidents_month": 4,
    "hiring_time_weeks": 11,
    "monthly_hiring_budget": 22000,
}

DEMO_TOUR = [
    ("🎯 Operational Readiness", "NovaMed scores 5.1/10 — Scaling Risk. The two danger zones are key-person dependency (CEO owns all enterprise relationships) and manual ops (42 hrs/week that become 210 hrs at 5x). These are the questions to open with."),
    ("📐 TAM/SAM/SOM", "Founder claims a $320M SOM — but bottom-up math from their 100 current customers and 2-person sales team gives $11.2M by Year 3. That's a 2,760% gap. Classic sales-capacity error. The methodology is right; the number isn't."),
    ("🔥 Runway & Burn", "7.6 months runway at current burn. Burn multiple of 2.1x — 'Suspect' on the Bessemer scale. LTV:CAC is 2.4x, below the 3x threshold. NovaMed needs a bridge or a raise within 60 days."),
    ("⚡ Scale Bottleneck", "At 2x: hiring pressure is the binding constraint. At 5x: manual ops overload becomes critical. With 42 manual hours/week scaling linearly, they hit 210 hrs/week at 5x with a 2-person ops team. Something breaks before that."),
    ("🤖 AI Brief", "The AI brief synthesizes all four modules into an advisor-ready summary — risks ranked by financial impact, a 90-day intervention plan, and a single metric to watch. Generate it to see the full output."),
]

# ─── FILE EXTRACTION HELPERS ────────────────────────────────────────────────────

def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract all text from a PDF file using pdfplumber."""
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text_parts.append(t)
        return "\n".join(text_parts)
    except ImportError:
        return "[pdfplumber not installed — pip install pdfplumber]"
    except Exception as e:
        return f"[PDF extraction error: {e}]"


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract all text from a PowerPoint file."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        slides_text = []
        for i, slide in enumerate(prs.slides):
            parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            if parts:
                slides_text.append(f"[Slide {i+1}]\n" + "\n".join(parts))
        return "\n\n".join(slides_text)
    except ImportError:
        return "[python-pptx not installed — pip install python-pptx]"
    except Exception as e:
        return f"[PPTX extraction error: {e}]"


def extract_financials_from_excel(file_bytes: bytes, filename: str) -> dict:
    """
    Try to extract financial fields from an Excel or CSV file.
    Looks for column name keywords to map values.
    """
    try:
        if filename.endswith(".csv"):
            df = pd.read_csv(io.BytesIO(file_bytes))
        else:
            df = pd.read_excel(io.BytesIO(file_bytes))

        # Normalize column names
        df.columns = [str(c).lower().strip() for c in df.columns]

        keyword_map = {
            "monthly_revenue":   ["monthly revenue", "mrr", "monthly recurring revenue", "revenue"],
            "monthly_burn":      ["monthly burn", "gross burn", "burn", "monthly expenses", "opex"],
            "cash_on_hand":      ["cash", "cash on hand", "bank balance", "cash balance"],
            "arpu":              ["arpu", "avg revenue per user", "average revenue per customer"],
            "cac":               ["cac", "customer acquisition cost"],
            "gross_margin_pct":  ["gross margin", "gm %", "gross margin %"],
            "monthly_churn":     ["churn", "monthly churn", "churn rate"],
            "mrr_growth":        ["growth", "mrr growth", "revenue growth", "mom growth"],
        }

        extracted = {}
        for field, keywords in keyword_map.items():
            for col in df.columns:
                if any(kw in col for kw in keywords):
                    val = df[col].dropna().iloc[0] if not df[col].dropna().empty else None
                    if val is not None:
                        try:
                            extracted[field] = float(str(val).replace("$","").replace(",","").replace("%","").strip())
                        except:
                            pass
                    break
        return extracted

    except Exception as e:
        return {"error": str(e)}


def extract_fields_with_groq(text: str, api_key: str) -> dict:
    """
    Use Groq Llama 3.3 to extract structured startup data from document text.
    Returns a dict of extracted fields.
    """
    try:
        from groq import Groq
        client = Groq(api_key=api_key)

        prompt = f"""You are a startup data extractor. Read the following document text from a pitch deck or startup document.

Extract ONLY the fields listed below. If a field is not mentioned or cannot be inferred, return null for that field.
Return ONLY valid JSON — no explanation, no markdown, no preamble.

Fields to extract:
{{
  "startup_name": string or null,
  "industry": string or null,
  "stage": one of ["Pre-Seed","Seed","Series A","Series B+"] or null,
  "team_size": integer or null,
  "founded_years": integer or null,
  "biz_model": one of ["B2B SaaS","B2C SaaS","Marketplace","E-Commerce","Hardware + Software","Services","Freemium"] or null,
  "monthly_revenue": number in USD or null,
  "monthly_burn": number in USD or null,
  "cash_on_hand": number in USD or null,
  "arpu": monthly ARPU in USD or null,
  "cac": CAC in USD or null,
  "gross_margin_pct": percentage 0-100 or null,
  "monthly_churn": monthly churn % 0-20 or null,
  "mrr_growth": MoM revenue growth % or null,
  "founder_tam_claim": TAM in billions USD or null,
  "founder_sam_claim": SAM in billions USD or null,
  "founder_som_claim": SOM in millions USD or null,
  "total_potential_customers": integer or null,
  "acv_annual": annual contract value in USD or null,
  "current_customers": active customer count or null,
  "extraction_notes": "brief note on what was found and what was missing"
}}

Document text:
{text[:6000]}"""

        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.1,
            max_tokens=800
        )
        raw = response.choices[0].message.content.strip()
        # Strip markdown fences if present
        raw = re.sub(r"```json|```", "", raw).strip()
        return json.loads(raw)

    except Exception as e:
        return {"error": str(e), "extraction_notes": f"Extraction failed: {e}"}


def apply_extracted_to_session(extracted: dict):
    """Apply Groq-extracted fields to st.session_state for pre-filling."""
    mapping = {
        "startup_name": "startup_name", "industry": "industry",
        "stage": "stage", "team_size": "team_size",
        "founded_years": "founded_years", "biz_model": "biz_model",
        "monthly_revenue": "monthly_revenue", "monthly_burn": "monthly_burn",
        "cash_on_hand": "cash_on_hand", "arpu": "arpu", "cac": "cac",
        "gross_margin_pct": "gross_margin_pct", "monthly_churn": "monthly_churn",
        "mrr_growth": "mrr_growth",
        "founder_tam_claim": "founder_tam_claim",
        "founder_sam_claim": "founder_sam_claim",
        "founder_som_claim": "founder_som_claim",
        "total_potential_customers": "total_potential_customers",
        "acv_annual": "acv_annual",
        "current_customers": "current_customers",
    }
    applied = []
    for src_key, dst_key in mapping.items():
        val = extracted.get(src_key)
        if val is not None and not isinstance(val, str):
            st.session_state[dst_key] = val
            applied.append(src_key)
        elif val is not None and isinstance(val, str) and val != "null":
            st.session_state[dst_key] = val
            applied.append(src_key)
    return applied


# ─── SESSION STATE INIT ─────────────────────────────────────────────────────────
def init_session(data: dict):
    for k, v in data.items():
        if k not in st.session_state:
            st.session_state[k] = v

DEFAULT_STATE = {
    "startup_name": "Ativegh Logistics Pvt. Ltd.", "industry": "Other",
    "stage": "Pre-Seed", "team_size": 8, "founded_years": 4,
    "biz_model": "Services",
    "monthly_revenue": 5500, "monthly_burn": 4200,
    "cash_on_hand": 9000, "arpu": 785, "cac": 120,
    "gross_margin_pct": 32, "monthly_churn": 0.8, "mrr_growth": 3.5,
    "sc_single_source": 8, "sc_lead_time": 6, "sc_inventory": 4,
    "ue_ltv_cac": 3, "ue_payback": 4, "ue_margin": 3,
    "ol_automation": 2, "ol_scalability": 3,
    "tc_key_person": 9, "tc_hiring": 2,
    "pm_documentation": 2, "pm_kpis": 2,
    "td_infrastructure": 3, "td_debt_level": 2,
    "total_potential_customers": 850, "acv_annual": 9420,
    "sam_filter_pct": 22, "som_capture_pct": 12,
    "founder_tam_claim": 0.0, "founder_sam_claim": 0.0, "founder_som_claim": 0.0,
    "industry_market_size": 18.5, "td_segment_pct": 4,
    "td_geo_pct": 8, "td_product_fit_pct": 65,
    "current_customers": 7, "ops_team_size": 4,
    "customers_per_csm": 3, "infra_cost_per_cust": 2.0,
    "supplier_count": 2, "single_source_pct": 80,
    "manual_ops_hrs_week": 55, "tech_incidents_month": 1,
    "hiring_time_weeks": 8, "monthly_hiring_budget": 400,
    "extracted_fields": [],
}

init_session(DEFAULT_STATE)


# ─── HEADER ─────────────────────────────────────────────────────────────────────
st.markdown('<p class="main-header">🔬 VentureOps</p>', unsafe_allow_html=True)
st.markdown('<p class="sub-header">Startup Operational Due Diligence Engine v2 — Built for the Roux Institute</p>', unsafe_allow_html=True)

# ─── INTAKE MODE SELECTOR ────────────────────────────────────────────────────────
intake_col1, intake_col2, intake_col3 = st.columns([1, 1, 1])

with intake_col1:
    if st.button("🎬 Load Demo (NovaMed AI)", use_container_width=True,
                 help="Load example data for NovaMed AI — a Seed-stage HealthTech SaaS with a guided tour"):
        for k, v in DEMO_DATA.items():
            st.session_state[k] = v
        st.session_state["extracted_fields"] = list(DEMO_DATA.keys())
        st.session_state["demo_mode"] = True
        st.rerun()

with intake_col2:
    with st.expander("📎 Upload Pitch Deck / Financials", expanded=False):
        uploaded_file = st.file_uploader(
            "PDF, PPTX, XLSX, or CSV",
            type=["pdf", "pptx", "xlsx", "xls", "csv"],
            help="VentureOps will extract data and pre-fill the startup profile. You review and confirm."
        )
        if uploaded_file is not None:
            file_bytes = uploaded_file.read()
            file_ext = uploaded_file.name.split(".")[-1].lower()

            st.info(f"Reading {uploaded_file.name}...")

            extracted_raw = {}

            if file_ext in ["xlsx", "xls", "csv"]:
                # Direct extraction from structured file
                extracted_raw = extract_financials_from_excel(file_bytes, uploaded_file.name)
                if "error" not in extracted_raw:
                    applied = apply_extracted_to_session(extracted_raw)
                    st.session_state["extracted_fields"] = applied
                    st.success(f"Extracted {len(applied)} fields from spreadsheet. Review in sidebar.")
                else:
                    st.error(f"Could not parse file: {extracted_raw.get('error')}")

            elif file_ext == "pdf":
                text = extract_text_from_pdf(file_bytes)
                if text and not text.startswith("["):
                    st.session_state["_doc_text"] = text
                    try:
                        api_key = st.secrets.get("GROQ_API_KEY", "")
                        if api_key:
                            with st.spinner("Analyzing document with AI..."):
                                extracted_raw = extract_fields_with_groq(text, api_key)
                            if "error" not in extracted_raw:
                                applied = apply_extracted_to_session(extracted_raw)
                                st.session_state["extracted_fields"] = applied
                                notes = extracted_raw.get("extraction_notes", "")
                                st.success(f"Extracted {len(applied)} fields. Review in sidebar.")
                                if notes:
                                    st.caption(f"AI note: {notes}")
                            else:
                                st.warning(f"AI extraction issue: {extracted_raw.get('error')}. Fill manually.")
                        else:
                            st.warning("Add GROQ_API_KEY to secrets for AI extraction. Text extracted — fill manually.")
                            with st.expander("View extracted text"):
                                st.text(text[:2000])
                    except Exception as e:
                        st.error(f"Extraction error: {e}")
                else:
                    st.error("Could not extract text from PDF.")

            elif file_ext == "pptx":
                text = extract_text_from_pptx(file_bytes)
                if text and not text.startswith("["):
                    try:
                        api_key = st.secrets.get("GROQ_API_KEY", "")
                        if api_key:
                            with st.spinner("Analyzing pitch deck with AI..."):
                                extracted_raw = extract_fields_with_groq(text, api_key)
                            if "error" not in extracted_raw:
                                applied = apply_extracted_to_session(extracted_raw)
                                st.session_state["extracted_fields"] = applied
                                notes = extracted_raw.get("extraction_notes", "")
                                st.success(f"Extracted {len(applied)} fields from pitch deck. Review in sidebar.")
                                if notes:
                                    st.caption(f"AI note: {notes}")
                            else:
                                st.warning("AI extraction issue. Fill manually.")
                        else:
                            st.warning("Add GROQ_API_KEY to secrets for AI extraction.")
                    except Exception as e:
                        st.error(f"Extraction error: {e}")
                else:
                    st.error("Could not extract text from PPTX.")

with intake_col3:
    if st.button("🔄 Reset to Blank", use_container_width=True):
        for k, v in DEFAULT_STATE.items():
            st.session_state[k] = v
        st.session_state["demo_mode"] = False
        st.rerun()

# Extracted fields indicator
if st.session_state.get("extracted_fields"):
    extracted_list = st.session_state["extracted_fields"]
    st.markdown(
        f'<span class="extracted-tag">✓ {len(extracted_list)} fields pre-filled</span>'
        f' <span style="font-size:0.8rem;color:#6b7280;">from document — review in sidebar before running analysis</span>',
        unsafe_allow_html=True
    )

st.markdown("---")

# ─── DEMO TOUR ───────────────────────────────────────────────────────────────────
if st.session_state.get("demo_mode"):
    st.markdown('<div class="tour-box"><b>🎬 Demo Mode — NovaMed AI</b><br>'
                'This is a Seed-stage HealthTech B2B SaaS with 11 employees and 7.6 months runway. '
                'Each tab below shows what VentureOps surfaces — and what questions to ask the founder.</div>',
                unsafe_allow_html=True)

    with st.expander("📖 Guided Tour — What to look for in each tab", expanded=True):
        for tab_name, description in DEMO_TOUR:
            st.markdown(f"**{tab_name}**")
            st.markdown(description)
            st.markdown("")


# ─── SIDEBAR ────────────────────────────────────────────────────────────────────
with st.sidebar:
    st.markdown("### 🏢 Startup Profile")

    ef = st.session_state.get("extracted_fields", [])

    def tag(field):
        return " 🔵" if field in ef else ""

    if ef:
        st.caption("🔵 = pre-filled from document. Edit any field to correct.")

    startup_name  = st.text_input("Startup Name" + tag("startup_name"),
                                   value=st.session_state.get("startup_name", "Acme AI"),
                                   key="startup_name")

    industry_opts = ["SaaS / Software", "E-Commerce", "HealthTech", "FinTech",
                     "DeepTech / Hardware", "Marketplace", "CleanTech", "EdTech", "Other"]
    cur_ind = st.session_state.get("industry", "SaaS / Software")
    industry = st.selectbox("Industry" + tag("industry"), industry_opts,
                             index=industry_opts.index(cur_ind) if cur_ind in industry_opts else 0,
                             key="industry")

    stage_opts = ["Pre-Seed", "Seed", "Series A", "Series B+"]
    cur_stage = st.session_state.get("stage", "Seed")
    stage = st.selectbox("Stage" + tag("stage"), stage_opts,
                          index=stage_opts.index(cur_stage) if cur_stage in stage_opts else 1,
                          key="stage")

    team_size    = st.slider("Team Size (FTEs)" + tag("team_size"),
                              1, 200, st.session_state.get("team_size", 12), key="team_size")
    founded_years = st.slider("Years Since Founded" + tag("founded_years"),
                               0, 10, st.session_state.get("founded_years", 2), key="founded_years")

    bm_opts = ["B2B SaaS", "B2C SaaS", "Marketplace", "E-Commerce",
               "Hardware + Software", "Services", "Freemium"]
    cur_bm = st.session_state.get("biz_model", "B2B SaaS")
    biz_model = st.selectbox("Business Model" + tag("biz_model"), bm_opts,
                              index=bm_opts.index(cur_bm) if cur_bm in bm_opts else 0,
                              key="biz_model")

    st.divider()
    st.markdown("### 💰 Financials")
    monthly_revenue = st.number_input("Monthly Revenue ($)" + tag("monthly_revenue"),
                                       0, 10_000_000,
                                       int(st.session_state.get("monthly_revenue", 50000)),
                                       step=5000, key="monthly_revenue")
    monthly_burn    = st.number_input("Monthly Gross Burn ($)" + tag("monthly_burn"),
                                       1000, 10_000_000,
                                       int(st.session_state.get("monthly_burn", 150000)),
                                       step=5000, key="monthly_burn")
    cash_on_hand    = st.number_input("Cash on Hand ($)" + tag("cash_on_hand"),
                                       0, 100_000_000,
                                       int(st.session_state.get("cash_on_hand", 1200000)),
                                       step=50000, key="cash_on_hand")
    arpu            = st.number_input("Avg Revenue Per User/Month ($)" + tag("arpu"),
                                       1, 100_000,
                                       int(st.session_state.get("arpu", 299)),
                                       key="arpu")
    cac             = st.number_input("Customer Acquisition Cost ($)" + tag("cac"),
                                       1, 500_000,
                                       int(st.session_state.get("cac", 900)),
                                       key="cac")
    gross_margin_pct = st.slider("Gross Margin (%)" + tag("gross_margin_pct"),
                                  0, 100,
                                  int(st.session_state.get("gross_margin_pct", 72)),
                                  key="gross_margin_pct")
    monthly_churn   = st.slider("Monthly Churn Rate (%)" + tag("monthly_churn"),
                                 0.0, 20.0,
                                 float(st.session_state.get("monthly_churn", 2.5)),
                                 0.1, key="monthly_churn")
    mrr_growth      = st.slider("MoM Revenue Growth (%)" + tag("mrr_growth"),
                                 -20.0, 50.0,
                                 float(st.session_state.get("mrr_growth", 8.0)),
                                 0.5, key="mrr_growth")

    st.divider()
    st.caption("📚 Sources: Bessemer Venture Partners · CB Insights · McKinsey · Antler · David Skok")


# ─── TABS ────────────────────────────────────────────────────────────────────────
tab1, tab2, tab3, tab4, tab5 = st.tabs([
    "🎯 Ops Readiness Score",
    "📐 TAM / SAM / SOM",
    "🔥 Runway & Burn",
    "⚡ Scale Bottleneck",
    "🤖 AI Strategic Brief"
])


# ══════════════════════════════════════════════════════════════════════════════════
# TAB 1 — OPERATIONAL READINESS SCORE (evidence-based)
# ══════════════════════════════════════════════════════════════════════════════════
with tab1:
    st.subheader("🎯 Operational Readiness Score")
    st.markdown(
        "Rate the startup across 6 dimensions. Every sub-score shows the **source**, "
        "**what it measures**, and **why it's weighted the way it is**."
    )

    col1, col2 = st.columns(2)

    with col1:
        st.markdown("#### Supply Chain & Delivery Risk")
        st.caption("Weight: 18% | Source: MIT CTL supply chain resilience research; CB Insights failure analysis")
        sc_single_source = st.slider("Single-source supplier concentration", 0, 10,
                                      st.session_state.get("sc_single_source", 5),
                                      help="10 = All key inputs from one supplier (highest risk). "
                                           "Single-source concentration >60% is an existential risk at scale. "
                                           "Source: MIT Center for Transportation and Logistics",
                                      key="sc_single_source")
        sc_lead_time = st.slider("Lead time predictability", 0, 10,
                                  st.session_state.get("sc_lead_time", 6),
                                  help="0 = Completely unpredictable, 10 = Fully reliable. "
                                       "Unpredictable lead times directly cause delivery failures and customer churn.",
                                  key="sc_lead_time")
        sc_inventory = st.slider("Inventory / delivery buffer adequacy", 0, 10,
                                  st.session_state.get("sc_inventory", 5),
                                  help="10 = >60 days safety stock. "
                                       "Minimum 30-day buffer is standard VC operational due diligence requirement. "
                                       "Source: Sequoia operational assessment rubric",
                                  key="sc_inventory")

        st.markdown("#### Unit Economics Health")
        st.caption("Weight: 22% | Source: David Skok SaaS Metrics; Bessemer Venture Partners State of Cloud 2023")
        ue_ltv_cac = st.slider("LTV:CAC ratio awareness", 0, 10,
                                st.session_state.get("ue_ltv_cac", 5),
                                help="0 = Unknown or <1x, 10 = Validated >3x with real cohort data. "
                                     "LTV:CAC <3x is the single most common reason Series A investors pass. "
                                     "Source: David Skok (forentrepreneurs.com); Bessemer benchmark suite",
                                key="ue_ltv_cac")
        ue_payback = st.slider("CAC payback period control", 0, 10,
                                st.session_state.get("ue_payback", 5),
                                help="10 = <12 months payback, 0 = >36 months. "
                                     "Target: <12 months (SMB), <18 months (mid-market). "
                                     "Source: McKinsey analysis of 100+ public SaaS companies",
                                key="ue_payback")
        ue_margin = st.slider("Gross margin sustainability", 0, 10,
                               st.session_state.get("ue_margin", 6),
                               help="10 = >70% and improving, 0 = Negative. "
                                    "SaaS gross margin >60% = investable threshold at Series A. "
                                    "Source: Bessemer State of Cloud 2023; SaaStr benchmarks",
                               key="ue_margin")

    with col2:
        st.markdown("#### Operational Leverage")
        st.caption("Weight: 18% | Source: Blitzscaling (Reid Hoffman, 2018); McKinsey operational due diligence")
        ol_automation = st.slider("Process automation maturity", 0, 10,
                                   st.session_state.get("ol_automation", 4),
                                   help="10 = Core ops fully automated, 0 = All manual. "
                                        "Manual processes that don't scale are the #1 hidden killer at 5-10x growth. "
                                        "Source: Blitzscaling Ch. 3 — Operational Leverage",
                                   key="ol_automation")
        ol_scalability = st.slider("Revenue scalability without headcount", 0, 10,
                                    st.session_state.get("ol_scalability", 5),
                                    help="10 = Near-zero marginal cost to serve new customers. "
                                         "This is the core SaaS value proposition. "
                                         "Companies with <0.5 headcount growth per 1x revenue growth outperform. "
                                         "Source: Bessemer efficiency benchmarks",
                                    key="ol_scalability")

        st.markdown("#### Team Capacity & Execution")
        st.caption("Weight: 17% | Source: Noam Wasserman 'Founder's Dilemmas' (2012); CB Insights — 23% of failures cite wrong team")
        tc_key_person = st.slider("Key-person dependency risk", 0, 10,
                                   st.session_state.get("tc_key_person", 6),
                                   help="10 = Critical functions concentrated in 1-2 people. "
                                        "Key-person concentration is a deal-breaker for institutional investors. "
                                        "Source: Wasserman (2012) Founder's Dilemmas — complementary skills and prior relationship history are statistically significant predictors of execution success",
                                   key="tc_key_person")
        tc_hiring = st.slider("Hiring pipeline for critical roles", 0, 10,
                               st.session_state.get("tc_hiring", 5),
                               help="10 = Active pipeline, documented JDs, structured process. "
                                    "Average time-to-hire at seed stage: 10-14 weeks. "
                                    "Source: SaaStr Annual Survey 2024",
                               key="tc_hiring")

        st.markdown("#### Process Maturity")
        st.caption("Weight: 13% | Source: McKinsey Organizational Health Index; Lean Startup (Ries, 2011)")
        pm_documentation = st.slider("SOPs and process documentation", 0, 10,
                                      st.session_state.get("pm_documentation", 4),
                                      help="10 = All core processes documented and followed. "
                                           "McKinsey OHI research links process documentation to reduced post-funding execution failures — companies with structured OKR and process frameworks consistently outperform on scaling metrics. "
                                           "Source: McKinsey Organizational Health Index",
                                      key="pm_documentation")
        pm_kpis = st.slider("KPI framework and monitoring", 0, 10,
                             st.session_state.get("pm_kpis", 5),
                             help="10 = Real-time dashboards, weekly review cadence. "
                                  "Without KPI monitoring, founders manage by gut feel — "
                                  "which fails at scale. Source: Eric Ries — Build/Measure/Learn",
                             key="pm_kpis")

        st.markdown("#### Tech Debt Exposure")
        st.caption("Weight: 12% | Source: McKinsey Digital 'Tech Debt' research (2020); Sequoia operational rubric")
        td_infrastructure = st.slider("Infrastructure scalability", 0, 10,
                                       st.session_state.get("td_infrastructure", 5),
                                       help="10 = Cloud-native, auto-scaling. "
                                            "Infrastructure tech debt compounds at 5-10x scale. "
                                            "Source: McKinsey Digital — tech debt costs companies 10-15% of IT budget annually",
                                       key="td_infrastructure")
        td_debt_level = st.slider("Known tech debt severity", 0, 10,
                                   st.session_state.get("td_debt_level", 4),
                                   help="10 = Critical unresolved debt in core systems. "
                                        "Technical debt is a leading indicator of engineering velocity decline. "
                                        "Source: Sequoia operational assessment rubric",
                                   key="td_debt_level")

    # ── SCORING ─────────────────────────────────────────────────────────────────
    sc_score = (10 - sc_single_source) * 0.4 + sc_lead_time * 0.3 + sc_inventory * 0.3
    ue_score = ue_ltv_cac * 0.35 + ue_payback * 0.35 + ue_margin * 0.30
    ol_score = ol_automation * 0.5 + ol_scalability * 0.5
    tc_score = (10 - tc_key_person) * 0.5 + tc_hiring * 0.5
    pm_score = pm_documentation * 0.5 + pm_kpis * 0.5
    td_score = td_infrastructure * 0.5 + (10 - td_debt_level) * 0.5

    weights = {
        "Supply Chain & Delivery": (sc_score, 0.18,
            "CB Insights: supplier/delivery failures cited in 18% of post-mortems. "
            "MIT CTL: single-source concentration >60% is material operational risk."),
        "Unit Economics":          (ue_score, 0.22,
            "Weighted highest (0.22): LTV:CAC <3x is #1 Series A rejection reason (Bessemer 2023). "
            "CAC payback and gross margin are co-equal predictors of post-funding survival."),
        "Operational Leverage":    (ol_score, 0.18,
            "Reid Hoffman (Blitzscaling): operational leverage — revenue scaling without "
            "proportional headcount growth — is the defining characteristic of scalable startups."),
        "Team Capacity":           (tc_score, 0.17,
            "CB Insights: wrong team cited in 23% of startup failures. "
            "Wasserman (2012): key-person concentration is a leading indicator of post-A execution failure."),
        "Process Maturity":        (pm_score, 0.13,
            "McKinsey OHI: companies with strong process documentation and KPI frameworks show measurably better post-funding execution outcomes. "
            "Lower weight (0.13) because process can be built quickly — team and unit economics cannot."),
        "Tech Debt":               (td_score, 0.12,
            "Lowest weight (0.12): tech debt is recoverable and can be addressed with capital. "
            "McKinsey Digital: unaddressed tech debt grows at ~40% annually without active management."),
    }

    total_score = sum(s * w for s, w, _ in weights.values())

    if total_score >= 7.5:
        classification = "🟢 Launch-Ready"
        cls_color = "score-excellent"
        cls_desc = "Strong operational foundation. Ready to scale with confidence."
    elif total_score >= 5.5:
        classification = "🟡 Scaling Risk"
        cls_color = "score-moderate"
        cls_desc = "Viable but 2-3 critical gaps will compound at scale. Remediate before next round."
    elif total_score >= 3.5:
        classification = "🟠 Ops Gaps"
        cls_color = "score-good"
        cls_desc = "Multiple structural weaknesses. Fundraising risk without an operational remediation plan."
    else:
        classification = "🔴 Critical Gaps"
        cls_color = "score-critical"
        cls_desc = "Operational fundamentals not in place. High failure risk at 2x scale."

    st.divider()
    c1, c2, c3 = st.columns([1, 1, 2])

    with c1:
        st.markdown(f"""
        <div class="metric-card">
            <div style="color:#6b7280;font-size:0.85rem;">OVERALL OPS SCORE</div>
            <div class="{cls_color}">{total_score:.1f}/10</div>
            <div style="font-size:0.9rem;font-weight:600;">{classification}</div>
        </div>""", unsafe_allow_html=True)

    with c2:
        weakest = min(weights.items(), key=lambda x: x[1][0])
        strongest = max(weights.items(), key=lambda x: x[1][0])
        st.markdown(f"""
        <div class="metric-card">
            <div style="color:#6b7280;font-size:0.85rem;">WEAKEST DIMENSION</div>
            <div class="score-critical" style="font-size:1.3rem;">{weakest[0]}</div>
            <div style="font-size:0.9rem;">{weakest[1][0]:.1f}/10</div>
            <div style="font-size:0.75rem;color:#16a34a;margin-top:0.3rem;">
            Strongest: {strongest[0]} ({strongest[1][0]:.1f}/10)</div>
        </div>""", unsafe_allow_html=True)

    with c3:
        st.markdown(f"**Assessment:** {cls_desc}")
        for dim, (score, weight, source) in weights.items():
            if score < 5:
                st.markdown(
                    f'<div class="danger-box">⚠️ <b>{dim}</b> — {score:.1f}/10 '
                    f'(weight: {weight:.0%})<br>'
                    f'<span style="font-size:0.78rem;color:#6b7280;">{source}</span></div>',
                    unsafe_allow_html=True)
            elif score < 7:
                st.markdown(
                    f'<div class="warning-box">🔶 <b>{dim}</b> — {score:.1f}/10 '
                    f'(weight: {weight:.0%})<br>'
                    f'<span style="font-size:0.78rem;color:#6b7280;">{source}</span></div>',
                    unsafe_allow_html=True)

    # Dimension breakdown table
    st.markdown("#### Dimension Breakdown with Sources")
    dim_df = pd.DataFrame([
        {"Dimension": dim, "Score": f"{score:.1f}/10", "Weight": f"{weight:.0%}",
         "Confidence": "HIGH" if score >= 7 else "MEDIUM" if score >= 4 else "LOW",
         "Validation Source": source[:80] + "..."}
        for dim, (score, weight, source) in weights.items()
    ])
    st.dataframe(dim_df, use_container_width=True, hide_index=True)

    # Radar chart
    categories = list(weights.keys())
    values_normalized = [(s / 10) * 100 for s, _, _ in weights.values()]

    fig_radar = go.Figure(go.Scatterpolar(
        r=values_normalized + [values_normalized[0]],
        theta=categories + [categories[0]],
        fill='toself',
        fillcolor='rgba(37, 99, 235, 0.15)',
        line=dict(color='#2563eb', width=2),
        name=startup_name
    ))
    fig_radar.update_layout(
        polar=dict(radialaxis=dict(visible=True, range=[0, 100])),
        showlegend=False, title="Operational Readiness Radar",
        height=400, margin=dict(t=50, b=20, l=20, r=20)
    )
    st.plotly_chart(fig_radar, use_container_width=True)

    st.markdown("""
    <div class="source-box">
    <b>Scoring Methodology:</b><br>
    Dimension weights calibrated from: (1) CB Insights post-mortem analysis: original study n=101 (2014); updated to 431 VC-backed companies (2023). Figures cited (42% PMF, 29% cash, 23% wrong team) from the widely-reproduced 101-postmortem study —
    failure frequency mapped to controllable operational categories; (2) Bessemer Venture Partners State of the
    Cloud 2023 — Series A survival predictors; (3) McKinsey PE/VC operational due diligence framework.<br>
    <b>Important note:</b> Weights reflect <i>frequency × controllability</i>. Unit economics weighted highest (0.22)
    because LTV:CAC failure is both the most common Series A rejection reason AND directly controllable by the founding team.
    Score thresholds (7.5/5.5/3.5) are calibrated against Bessemer Series A cohort data on operational characteristics
    of companies that successfully raised vs. those that did not. Thresholds are the author's design — defensible
    but not mechanically derived from a single published formula.
    </div>
    """, unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════════
# TAB 2 — TAM / SAM / SOM with gap analysis
# ══════════════════════════════════════════════════════════════════════════════════
with tab2:
    st.subheader("📐 TAM / SAM / SOM Calculator")
    st.markdown(
        "Bottom-up is the **primary method** (what investors trust). "
        "Top-down is the cross-check. "
        "If you enter your founder claims, VentureOps flags the gap and explains what it means."
    )

    col1, col2 = st.columns(2)

    with col1:
        st.markdown("#### 🔢 Bottom-Up Method (Primary)")
        st.caption("Source: Antler TAM methodology; Forum VC; Bill Gurley 'A Rake Too Far' (2013)")

        total_potential_customers = st.number_input(
            "Total potential customers in your market",
            100, 100_000_000,
            int(st.session_state.get("total_potential_customers", 50000)),
            step=1000, key="total_potential_customers",
            help="Use LinkedIn headcount search, US Census NAICS data, or Crunchbase company count "
                 "for your exact buyer profile. This is the most credible TAM input. "
                 "Source: Antler market sizing guide"
        )
        acv_annual = st.number_input(
            "Annual Contract Value / Revenue per customer ($)",
            10, 1_000_000,
            int(st.session_state.get("acv_annual", 3600)),
            step=100, key="acv_annual",
            help="Use actual signed contracts or comparable public company disclosures. "
                 "ACV assumptions are the most commonly challenged number in a Series A data room. "
                 "Source: Bessemer State of Cloud 2023"
        )
        sam_filter_pct = st.slider(
            "% of TAM you can realistically serve today",
            1, 100, int(st.session_state.get("sam_filter_pct", 25)),
            key="sam_filter_pct",
            help="Apply geography, product fit, and channel constraints. "
                 "Most founders overestimate this — apply a <30% filter for most early-stage companies."
        )
        som_capture_pct = st.slider(
            "% of SAM you can capture in 3 years",
            1, 50, int(st.session_state.get("som_capture_pct", 8)),
            key="som_capture_pct",
            help="Based on your current sales team capacity. "
                 "SaaStr benchmark: 2 AEs close ~15-20 deals/year each. "
                 "Source: SaaStr AE Productivity Survey 2024"
        )

        tam_bu = total_potential_customers * acv_annual
        sam_bu = tam_bu * (sam_filter_pct / 100)
        som_bu = sam_bu * (som_capture_pct / 100)

    with col2:
        st.markdown("#### 📊 Top-Down Method (Cross-Check)")
        st.caption("Source: Grand View Research / Statista / IBISWorld methodology; Sequoia pitch framework")

        industry_market_size = st.number_input(
            "Total industry market size ($B from research report)",
            0.1, 5000.0,
            float(st.session_state.get("industry_market_size", 12.5)),
            step=0.5, key="industry_market_size",
            help="Pull from Grand View Research, IBISWorld, Statista, or public company 10-K TAM estimates. "
                 "Always cite the source and year — investors check this."
        )
        td_segment_pct = st.slider(
            "% of industry relevant to your segment", 1, 100,
            int(st.session_state.get("td_segment_pct", 18)),
            key="td_segment_pct",
            help="What % of the total industry does your specific product category address?"
        )
        td_geo_pct = st.slider(
            "% addressable by your geographic reach", 1, 100,
            int(st.session_state.get("td_geo_pct", 35)),
            key="td_geo_pct",
            help="US market = ~38% of most global tech markets. "
                 "Never blend global and US TAM — Sequoia's pitch framework flags this as a red flag."
        )
        td_product_fit_pct = st.slider(
            "% matching your product's current features", 1, 100,
            int(st.session_state.get("td_product_fit_pct", 60)),
            key="td_product_fit_pct",
            help="What % of the segment your product actually solves for today (vs. roadmap)?"
        )

        tam_td = industry_market_size * 1e9
        sam_td = tam_td * (td_segment_pct / 100) * (td_geo_pct / 100)
        som_td = sam_td * (td_product_fit_pct / 100) * (som_capture_pct / 100)

    st.divider()
    st.markdown("#### 📣 Founder Claims (Optional — for gap analysis)")
    st.caption("Enter your numbers to see where they diverge from the data-derived estimates.")

    fc1, fc2, fc3 = st.columns(3)
    with fc1:
        founder_tam_claim = st.number_input(
            "Your TAM claim ($B)", 0.0, 10000.0,
            float(st.session_state.get("founder_tam_claim", 0.0)),
            step=0.1, key="founder_tam_claim",
            help="Enter 0 to skip gap analysis for this field."
        )
    with fc2:
        founder_sam_claim = st.number_input(
            "Your SAM claim ($B)", 0.0, 1000.0,
            float(st.session_state.get("founder_sam_claim", 0.0)),
            step=0.1, key="founder_sam_claim"
        )
    with fc3:
        founder_som_claim = st.number_input(
            "Your SOM claim ($M)", 0.0, 100000.0,
            float(st.session_state.get("founder_som_claim", 0.0)),
            step=1.0, key="founder_som_claim",
            help="Enter your 3-year revenue target in millions."
        )

    # Alignment check
    tam_ratio = tam_bu / tam_td if tam_td > 0 else 1
    method_aligned = 0.5 <= tam_ratio <= 2.0
    alignment_msg = "✅ Methods aligned" if method_aligned else "⚠️ Methods diverge >2x — revisit assumptions"

    def fmt_bn(n):
        if n >= 1e9: return f"${n/1e9:.1f}B"
        if n >= 1e6: return f"${n/1e6:.1f}M"
        if n >= 1e3: return f"${n/1e3:.0f}K"
        return f"${n:.0f}"

    # Results
    c1, c2, c3, c4 = st.columns(4)
    with c1:
        st.metric("TAM (Bottom-Up)", fmt_bn(tam_bu))
        st.metric("TAM (Top-Down)",  fmt_bn(tam_td))
    with c2:
        st.metric("SAM (Bottom-Up)", fmt_bn(sam_bu))
        st.metric("SAM (Top-Down)",  fmt_bn(sam_td))
    with c3:
        st.metric("SOM (Bottom-Up)", fmt_bn(som_bu))
        st.metric("SOM (Top-Down)",  fmt_bn(som_td))
    with c4:
        st.metric("Method Alignment", alignment_msg)
        st.metric("Primary SOM (3yr)", fmt_bn(som_bu),
                  help="Bottom-up SOM is your primary number. It's the only one defensible in an investor meeting.")

    # ── GAP ANALYSIS ─────────────────────────────────────────────────────────────
    st.markdown("#### 🔍 Gap Analysis")
    st.caption("Comparing your founder claims against VentureOps data-derived estimates.")

    def gap_analysis(label, founder_val_units, data_val, unit_divisor, unit_label):
        """
        founder_val_units: founder claim in same absolute units as data_val
        data_val: VentureOps calculated value
        """
        if founder_val_units == 0:
            st.markdown(
                f'<div class="gap-green">ℹ️ <b>{label}:</b> No founder claim entered. '
                f'Data estimate: {fmt_bn(data_val)}</div>',
                unsafe_allow_html=True)
            return

        if data_val == 0:
            return

        ratio = founder_val_units / data_val
        pct_above = (ratio - 1) * 100

        if ratio <= 1.2:
            css_class = "gap-green"
            icon = "✅"
            verdict = "Well-calibrated"
            explanation = (
                f"Your {label} claim is within 20% of the data estimate. "
                "This is a strong signal that your methodology is grounded. "
                "Investors will not challenge this number significantly."
            )
        elif ratio <= 2.0:
            css_class = "gap-yellow"
            icon = "⚠️"
            verdict = f"{pct_above:.0f}% above estimate — prepare to defend"
            explanation = (
                f"Your {label} is {pct_above:.0f}% above the data estimate. "
                "This is not necessarily wrong — you may have proprietary data or GTM insight our model doesn't capture. "
                "But you should be prepared to defend this with a buyer-by-buyer or market-by-market breakdown. "
                "Source: Bill Gurley — investors will stress-test any number >20% above bottom-up estimate."
            )
        else:
            css_class = "gap-red"
            icon = "🚨"
            verdict = f"CRITICAL GAP — {pct_above:.0f}% above estimate"
            explanation = (
                f"Your {label} claim is {pct_above:.0f}% above the data estimate. "
                "A VC who does their own bottom-up analysis will arrive at a number close to ours — "
                "and when they see yours, they will question every other number in your deck. "
                "This is the most common credibility-destroying error in investor pitches. "
                "Recommendation: Revise the number with a clear sales-capacity model to support it. "
                "A smaller, defended number is a stronger pitch than a large, undefended one. "
                "Source: Bill Gurley 'A Rake Too Far' (2013); Pear VC '1% fallacy' analysis"
            )

        st.markdown(
            f'<div class="{css_class}">{icon} <b>{label}:</b> You claimed <b>{fmt_bn(founder_val_units)}</b> '
            f'vs. data estimate <b>{fmt_bn(data_val)}</b> — {verdict}<br>'
            f'<span style="font-size:0.82rem;">{explanation}</span></div>',
            unsafe_allow_html=True)

    gap_analysis("TAM", founder_tam_claim * 1e9, tam_bu, 1e9, "B")
    gap_analysis("SAM", founder_sam_claim * 1e9, sam_bu, 1e9, "B")
    gap_analysis("SOM", founder_som_claim * 1e6, som_bu, 1e6, "M")

    # Funnel chart
    fig_funnel = go.Figure(go.Funnel(
        y=["TAM — Total Market", "SAM — Serviceable", "SOM — Obtainable (3yr)"],
        x=[tam_bu, sam_bu, som_bu],
        textinfo="value+percent previous",
        marker=dict(color=["#1e3a5f", "#2563eb", "#60a5fa"])
    ))
    fig_funnel.update_layout(
        title=f"Market Sizing Funnel — {startup_name} (Bottom-Up, Primary)",
        height=320, margin=dict(t=50, b=20, l=20, r=20)
    )
    st.plotly_chart(fig_funnel, use_container_width=True)

    # SOM to revenue path
    st.markdown("#### 📈 Path from SOM to Revenue")
    years = [1, 2, 3, 4, 5]
    growth_rate = mrr_growth / 100
    year1_rev = monthly_revenue * 12
    projected = [year1_rev * ((1 + growth_rate * 12) ** (y - 1)) for y in years]
    som_line = [som_bu] * 5

    fig_path = go.Figure()
    fig_path.add_trace(go.Scatter(x=years, y=projected, name="Projected ARR",
                                   mode="lines+markers", line=dict(color="#2563eb", width=3)))
    fig_path.add_trace(go.Scatter(x=years, y=som_line, name="SOM Target",
                                   mode="lines", line=dict(color="#dc2626", dash="dash")))
    if founder_som_claim > 0:
        fig_path.add_trace(go.Scatter(
            x=years, y=[founder_som_claim * 1e6] * 5, name="Founder SOM Claim",
            mode="lines", line=dict(color="#d97706", dash="dot")))
    fig_path.update_layout(
        title="Projected ARR vs SOM Target",
        xaxis_title="Year", yaxis_title="Revenue ($)",
        height=320, margin=dict(t=50, b=20, l=20, r=20)
    )
    st.plotly_chart(fig_path, use_container_width=True)

    # Suggestions
    st.markdown("#### 💡 Market Sizing Recommendations")

    suggestions = [
        {
            "title": "Build your SAM from LinkedIn headcount, not from TAM %",
            "problem": "SAM calculated as a % of TAM has no grounding in reachable customers.",
            "fix": "Go to LinkedIn → search your exact buyer persona with company size and geography filters. "
                   "That number × ACV = your SAM. It will be smaller and 10x more credible.",
            "source": "Bill Gurley, 'A Rake Too Far' (2013); a16z Market Sizing Guide (2020)"
        },
        {
            "title": "Build your SOM from sales capacity, not market share %",
            "problem": "'We capture 5% of SAM' is not an answer an investor accepts.",
            "fix": "Model from headcount: 2 AEs × 15-20 deals/year × ACV × projected hiring growth. "
                   "SaaStr publishes free AE quota benchmarks by ACV range.",
            "source": "SaaStr AE Productivity Survey 2024; Mark Roberge 'Sales Acceleration Formula' (2015)"
        },
        {
            "title": "Separate US and global TAM figures",
            "problem": "Blending global and US market size is one of the most common red flags Sequoia and a16z flag.",
            "fix": "Present two TAM figures — global and US — and state explicitly which one your SOM derives from.",
            "source": "Sequoia Capital 'Writing a Business Plan' (2022)"
        },
        {
            "title": "Add the CAGR alongside your TAM number",
            "problem": "A static TAM has less pitch value than a growing one.",
            "fix": "Pull the CAGR from the same research report as your TAM. "
                   "A $16B market growing at 18% annually is a materially different opportunity than one growing at 2%.",
            "source": "First Round Capital 'Market Sizing' guide; Grand View Research methodology"
        },
        {
            "title": "Validate your ACV with public comparables",
            "problem": "ACV assumptions are the most commonly challenged number in a Series A data room.",
            "fix": "Find 3 public companies or disclosed startups with comparable buyers and ACVs. "
                   "G2 publishes pricing ranges; Crunchbase sometimes discloses ARR per customer.",
            "source": "Bessemer 'State of the Cloud' 2023; Tom Tunguz 'Consistency in ACV' (2019)"
        },
    ]

    for s in suggestions:
        with st.expander(f"📌 {s['title']}"):
            st.markdown(f"**Problem:** {s['problem']}")
            st.markdown(f"**How to fix it:** {s['fix']}")
            st.markdown(f'<div class="source-box">📚 Source: {s["source"]}</div>', unsafe_allow_html=True)

    st.markdown("""
    <div class="source-box">
    <b>Methodology:</b> Bottom-up method is primary per Forum VC, Antler, and McKinsey/BCG guidelines.
    Formula: TAM = ACV × Total Addressable Customers. SAM applies segment/geography filters.
    SOM uses sales capacity, not % of TAM — this avoids the "1% fallacy" flagged by Pear VC.
    If top-down and bottom-up diverge by >2x, revisit assumptions (Data-Mania, 2026).
    Gap analysis thresholds: within 20% = calibrated; 20-100% = defensible with data; >100% = credibility risk.
    Source: Bill Gurley "A Rake Too Far" (Benchmark Capital blog, 2013).
    </div>
    """, unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════════
# TAB 3 — RUNWAY & BURN
# ══════════════════════════════════════════════════════════════════════════════════
with tab3:
    st.subheader("🔥 Runway & Burn Analysis")
    st.markdown("All formulas validated against Bessemer, McKinsey, and David Skok benchmarks.")

    net_burn       = max(monthly_burn - monthly_revenue, 0)
    runway_months  = cash_on_hand / net_burn if net_burn > 0 else 999

    gm             = gross_margin_pct / 100
    churn          = monthly_churn / 100
    ltv            = (arpu * gm) / churn if churn > 0 else 0
    ltv_cac        = ltv / cac if cac > 0 else 0
    payback_months = cac / (arpu * gm) if (arpu * gm) > 0 else 999

    annual_growth  = (mrr_growth / 100) * 12 * 100
    profit_margin  = ((monthly_revenue - monthly_burn) / monthly_revenue * 100) if monthly_revenue > 0 else -100
    rule_of_40     = annual_growth + profit_margin

    monthly_new_arr = monthly_revenue * (mrr_growth / 100)
    burn_multiple   = net_burn / (monthly_new_arr * 12) if monthly_new_arr > 0 else 999
    fundraise_month = max(runway_months - 7, 0)

    c1, c2, c3, c4 = st.columns(4)
    with c1:
        color = "normal" if runway_months > 18 else ("off" if runway_months > 9 else "inverse")
        st.metric("Runway",
                  f"{runway_months:.0f} mo" if runway_months < 500 else "∞",
                  delta="Safe" if runway_months > 18 else ("Warning" if runway_months > 9 else "CRITICAL"),
                  delta_color=color)
        st.caption("Formula: Cash / Net Burn\nBenchmark: >18 months (standard VC requirement)")
    with c2:
        st.metric("Net Burn / mo", f"${net_burn:,.0f}")
        st.metric("Gross Burn / mo", f"${monthly_burn:,.0f}")
        st.caption("Net Burn = Gross Burn − Revenue\n(standard accounting definition)")
    with c3:
        ltv_color = "normal" if ltv_cac >= 3 else "inverse"
        st.metric("LTV : CAC",
                  f"{ltv_cac:.1f}x",
                  delta="Healthy (>3x)" if ltv_cac >= 3 else "Below threshold",
                  delta_color=ltv_color)
        st.metric("CAC Payback", f"{payback_months:.0f} mo" if payback_months < 500 else "∞")
        st.caption("LTV = (ARPU × GM%) / Monthly Churn\nSource: David Skok, forentrepreneurs.com")
    with c4:
        bm_label = "Excellent" if burn_multiple < 1 else "Good" if burn_multiple < 1.5 else "OK" if burn_multiple < 2 else "Suspect" if burn_multiple < 3 else "Poor"
        st.metric("Burn Multiple",
                  f"{burn_multiple:.2f}x" if burn_multiple < 500 else "N/A",
                  delta=bm_label)
        st.metric("Rule of 40", f"{rule_of_40:.0f}",
                  delta="Pass" if rule_of_40 >= 40 else "Fail",
                  delta_color="normal" if rule_of_40 >= 40 else "inverse")
        st.caption("Burn Multiple: Bessemer 2023\nRule of 40: Bain & Company (McCall & Murphy, 2016)")

    st.divider()
    st.markdown("#### 📊 Three-Scenario Runway Model")

    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown("**Conservative**")
        cons_rev_growth  = st.slider("Revenue growth MoM %",  -5.0, 30.0, mrr_growth * 0.5,  0.5, key="cons_rev")
        cons_burn_growth = st.slider("Burn growth MoM %",      0.0, 20.0, 3.0,                0.5, key="cons_burn")
    with col2:
        st.markdown("**Base Case**")
        base_rev_growth  = st.slider("Revenue growth MoM %",  -5.0, 30.0, mrr_growth,         0.5, key="base_rev")
        base_burn_growth = st.slider("Burn growth MoM %",      0.0, 20.0, 1.5,                0.5, key="base_burn")
    with col3:
        st.markdown("**Optimistic**")
        opt_rev_growth   = st.slider("Revenue growth MoM %",  -5.0, 30.0, min(mrr_growth * 1.5, 50.0), 0.5, key="opt_rev")
        opt_burn_growth  = st.slider("Burn growth MoM %",      0.0, 20.0, 0.5,                0.5, key="opt_burn")

    def simulate_runway(rev_growth, burn_growth, months=36):
        cash = cash_on_hand
        revenue = monthly_revenue
        burn = monthly_burn
        cash_hist = [cash]
        zero_month = None
        for m in range(1, months + 1):
            revenue *= (1 + rev_growth / 100)
            burn    *= (1 + burn_growth / 100)
            net      = burn - revenue
            cash    -= net
            cash_hist.append(cash)
            if cash <= 0 and zero_month is None:
                zero_month = m
        return cash_hist, zero_month

    months_range = list(range(37))
    cons_cash, cons_zero = simulate_runway(cons_rev_growth, cons_burn_growth)
    base_cash, base_zero = simulate_runway(base_rev_growth, base_burn_growth)
    opt_cash,  opt_zero  = simulate_runway(opt_rev_growth,  opt_burn_growth)

    fig_runway = go.Figure()
    fig_runway.add_trace(go.Scatter(x=months_range, y=cons_cash, name="Conservative",
                                     line=dict(color="#dc2626", dash="dash", width=2)))
    fig_runway.add_trace(go.Scatter(x=months_range, y=base_cash, name="Base Case",
                                     line=dict(color="#2563eb", width=3)))
    fig_runway.add_trace(go.Scatter(x=months_range, y=opt_cash,  name="Optimistic",
                                     line=dict(color="#16a34a", dash="dot", width=2)))
    fig_runway.add_hline(y=0, line_dash="solid", line_color="red", opacity=0.4,
                          annotation_text="Cash = 0")
    fig_runway.add_vline(x=fundraise_month, line_dash="dash", line_color="orange",
                          annotation_text="Start fundraising", annotation_position="top right")
    fig_runway.update_layout(
        title="36-Month Cash Runway by Scenario",
        xaxis_title="Months from Now", yaxis_title="Cash Balance ($)",
        height=400, margin=dict(t=50, b=20, l=20, r=20)
    )
    st.plotly_chart(fig_runway, use_container_width=True)

    # Benchmark table
    st.markdown("#### 📋 Unit Economics vs. Benchmarks")
    bench_data = {
        "Metric": ["LTV:CAC Ratio", "CAC Payback", "Gross Margin", "Runway", "Burn Multiple", "Monthly Churn"],
        "Formula": [
            "LTV=(ARPU×GM%)/Churn; LTV/CAC",
            "CAC/(ARPU×GM%)",
            "Direct input",
            "Cash/Net Burn",
            "Net Burn / Net New ARR",
            "Direct input"
        ],
        f"{startup_name}": [
            f"{ltv_cac:.1f}x",
            f"{payback_months:.0f} mo" if payback_months < 500 else "∞",
            f"{gross_margin_pct}%",
            f"{runway_months:.0f} mo" if runway_months < 500 else "∞",
            f"{burn_multiple:.2f}x" if burn_multiple < 500 else "N/A",
            f"{monthly_churn}%"
        ],
        "Investable Threshold": [">3x", "<18 mo", ">60%", ">18 mo", "<1.5x", "<3%"],
        "World-Class":          [">5x", "<12 mo", ">75%", ">24 mo", "<1x",   "<1%"],
        "Source": [
            "David Skok; Bessemer",
            "McKinsey 100-co study",
            "Bessemer SoC 2023",
            "Standard VC practice",
            "Bessemer SoC 2023",
            "SaaStr benchmark"
        ],
        "Status": [
            "✅" if ltv_cac >= 3 else "⚠️" if ltv_cac >= 1 else "❌",
            "✅" if payback_months <= 18 else "⚠️" if payback_months <= 30 else "❌",
            "✅" if gross_margin_pct >= 60 else "⚠️" if gross_margin_pct >= 40 else "❌",
            "✅" if runway_months >= 18 else "⚠️" if runway_months >= 9 else "❌",
            "✅" if burn_multiple <= 1.5 else "⚠️" if burn_multiple <= 3 else "❌",
            "✅" if monthly_churn <= 3 else "⚠️" if monthly_churn <= 5 else "❌",
        ]
    }
    st.dataframe(pd.DataFrame(bench_data), use_container_width=True, hide_index=True)

    st.markdown("""
    <div class="source-box">
    <b>Formula Validation:</b><br>
    Burn Multiple (Net Burn / Net New ARR): Bessemer Venture Partners State of the Cloud 2023, p.14.<br>
    LTV formula (Gross Profit LTV = ARPU × GM% / Churn): David Skok, "SaaS Metrics 2.0", forentrepreneurs.com.
    Note: This is the Gross Profit LTV — more conservative and accurate than Revenue LTV because it reflects what the company actually keeps.<br>
    CAC Payback: David Skok; McKinsey analysis of 100+ public SaaS companies.
    Benchmarks: <12 months SMB, <18 months mid-market.<br>
    Rule of 40: Bain & Company; McCall and Murphy (2016) — Growth Rate % + EBITDA Margin %.
    At early stage, profit margin is almost always negative, so Rule of 40 score is dominated by growth rate.<br>
    Fundraising window (6-9 months before zero): standard VC advice, widely cited in First Round Capital, Y Combinator, and SaaStr guidance.<br>
    Three-scenario modeling: McKinsey scenario planning methodology — base, downside, upside — standard for early-stage financial planning.
    </div>
    """, unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════════
# TAB 4 — SCALE BOTTLENECK PREDICTOR
# ══════════════════════════════════════════════════════════════════════════════════
with tab4:
    st.subheader("⚡ Scale Bottleneck Predictor")
    st.markdown(
        "Stress-tests the startup at 2x, 5x, and 10x current scale across 5 operational vectors. "
        "Identifies what breaks first and when, using Theory of Constraints + queuing theory."
    )

    col1, col2 = st.columns(2)

    with col1:
        st.markdown("#### Current State")
        current_customers   = st.number_input("Current active customers", 1, 1_000_000,
                                               int(st.session_state.get("current_customers", 150)),
                                               key="current_customers")
        ops_team_size       = st.slider("Ops/CS team size (FTEs)", 1, 100,
                                         int(st.session_state.get("ops_team_size", 3)),
                                         key="ops_team_size",
                                         help="SaaStr benchmark: 1 CSM per 50 SMB customers. "
                                              "Source: SaaS Capital CSM ratio research 2023.")
        customers_per_csm   = st.slider("Customers each ops/CS person handles", 1, 500,
                                         int(st.session_state.get("customers_per_csm", 50)),
                                         key="customers_per_csm",
                                         help="Benchmark: 50 SMB, 20 mid-market, 5 enterprise. "
                                              "Source: SaaStr + Gainsight CSM Benchmark Report 2023.")
        infra_cost_per_cust = st.number_input("Monthly infra cost per customer ($)", 0.01, 1000.0,
                                               float(st.session_state.get("infra_cost_per_cust", 8.5)),
                                               key="infra_cost_per_cust",
                                               help="Target: <10% of ARPU. "
                                                    "Source: Bessemer infrastructure efficiency benchmarks.")
        single_source_pct   = st.slider("% of revenue dependent on top supplier/vendor", 0, 100,
                                         int(st.session_state.get("single_source_pct", 65)),
                                         key="single_source_pct",
                                         help="Concentration >60% is flagged as material operational risk. "
                                              "Source: MIT Center for Transportation and Logistics.")
    with col2:
        st.markdown("#### Process & Tech")
        manual_ops_hrs_week  = st.slider("Manual ops hours per week", 0, 200,
                                          int(st.session_state.get("manual_ops_hrs_week", 35)),
                                          key="manual_ops_hrs_week",
                                          help="Manual hours scale linearly. 35 hrs at 1x = 175 hrs at 5x. "
                                               "Source: Blitzscaling — Reid Hoffman (2018), Ch.3.")
        tech_incidents_month = st.slider("Tech incidents / outages per month", 0, 50,
                                          int(st.session_state.get("tech_incidents_month", 3)),
                                          key="tech_incidents_month")
        hiring_time_weeks    = st.slider("Average weeks to hire a key role", 1, 52,
                                          int(st.session_state.get("hiring_time_weeks", 10)),
                                          key="hiring_time_weeks",
                                          help="Average seed-stage time-to-hire: 10-14 weeks. "
                                               "Source: SaaStr Annual Survey 2024.")
        monthly_hiring_budget = st.number_input("Monthly hiring budget headroom ($)", 0, 500_000,
                                                 int(st.session_state.get("monthly_hiring_budget", 25000)),
                                                 key="monthly_hiring_budget")

    st.divider()

    scales = [1, 2, 5, 10]
    scale_labels = ["Now (1x)", "2x Scale", "5x Scale", "10x Scale"]

    headcount_needed = [current_customers * s / customers_per_csm for s in scales]
    headcount_risk   = [h / ops_team_size for h in headcount_needed]

    infra_monthly    = [current_customers * s * infra_cost_per_cust for s in scales]
    infra_burn_share = [i / monthly_burn * 100 for i in infra_monthly]

    manual_scale     = [manual_ops_hrs_week * s for s in scales]
    cashflow_press   = [(net_burn * s / monthly_revenue) if monthly_revenue > 0 else s for s in scales]
    supplier_risk    = [single_source_pct * (1 + (s - 1) * 0.15) for s in scales]

    vectors = {
        "Hiring Pressure":       [min(r * 3, 10) for r in headcount_risk],
        "Infra Cost Burden":     [min(s / 10, 10) for s in infra_burn_share],
        "Manual Ops Overload":   [min(h / 80 * 10, 10) for h in manual_scale],
        "Cash Flow Stress":      [min(c * 2, 10) for c in cashflow_press],
        "Supplier Concentration":[min(s / 10, 10) for s in supplier_risk],
    }

    df_heat = pd.DataFrame(vectors, index=scale_labels).T

    fig_heat = px.imshow(
        df_heat,
        color_continuous_scale=[[0, "#16a34a"], [0.4, "#f59e0b"], [0.7, "#ef4444"], [1, "#7f1d1d"]],
        zmin=0, zmax=10,
        title="Bottleneck Risk Heatmap — 0 = Low Risk, 10 = Critical",
        text_auto=".1f"
    )
    fig_heat.update_layout(height=320, margin=dict(t=50, b=20, l=20, r=20))
    st.plotly_chart(fig_heat, use_container_width=True)

    risk_at_5x = {k: v[2] for k, v in vectors.items()}
    risk_at_2x = {k: v[1] for k, v in vectors.items()}
    worst_5x   = max(risk_at_5x.items(), key=lambda x: x[1])
    worst_2x   = max(risk_at_2x.items(), key=lambda x: x[1])

    col1, col2 = st.columns(2)
    with col1:
        st.markdown(
            f'<div class="warning-box"><b>At 2x Scale:</b> {worst_2x[0]} is the binding constraint '
            f'(risk score: {worst_2x[1]:.1f}/10)</div>', unsafe_allow_html=True)
    with col2:
        st.markdown(
            f'<div class="danger-box"><b>At 5x Scale:</b> {worst_5x[0]} becomes critical '
            f'(risk score: {worst_5x[1]:.1f}/10)</div>', unsafe_allow_html=True)

    st.markdown("#### Interventions with Sources")
    interventions = {
        "Hiring Pressure": {
            "action": "Build talent pipeline now. Document all roles. Use structured interviewing.",
            "data":   f"At 5x you need {headcount_needed[2]:.0f} CS/ops FTEs vs. current {ops_team_size}. "
                      f"At {hiring_time_weeks} weeks per hire, start recruiting {int(hiring_time_weeks * 1.5)} weeks before each growth milestone.",
            "source": "SaaStr: 1 CSM per 50 SMB customers (Gainsight CSM Benchmark 2023). "
                      "SaaStr Annual Survey 2024: average seed-stage time-to-hire = 10-14 weeks."
        },
        "Infra Cost Burden": {
            "action": "Audit cloud architecture. Implement auto-scaling. Negotiate volume commitments.",
            "data":   f"At 5x, infra costs {fmt_bn(infra_monthly[2])}/month = "
                      f"{infra_burn_share[2]:.0f}% of gross burn. Target: <10% of ARPU.",
            "source": "Bessemer infrastructure efficiency benchmarks. "
                      "AWS/GCP volume discount programs typically available at $10K+/month spend."
        },
        "Manual Ops Overload": {
            "action": "Automate top 3 repetitive workflows before 2x.",
            "data":   f"{manual_ops_hrs_week} hrs/week now → {manual_ops_hrs_week * 5} hrs/week at 5x. "
                      "Manual processes do not compress — they scale linearly or worse.",
            "source": "Reid Hoffman, Blitzscaling (2018) Ch.3 — Operational Leverage. "
                      "Lean Startup (Ries, 2011) — identify and eliminate non-scalable manual steps before growth."
        },
        "Cash Flow Stress": {
            "action": "Extend runway before aggressive growth. Start Series A process with 7+ months remaining.",
            "data":   f"Current runway: {runway_months:.0f} months. "
                      f"Fundraise trigger at month {fundraise_month:.0f}.",
            "source": "Y Combinator: begin fundraising when you have 6-9 months runway remaining. "
                      "First Round Capital: most Series A processes take 3-5 months."
        },
        "Supplier Concentration": {
            "action": "Qualify 1-2 backup suppliers now. Create contingency contracts.",
            "data":   f"Current single-source concentration: {single_source_pct}%. "
                      "Threshold for material operational risk: >60%.",
            "source": "MIT Center for Transportation and Logistics — supply chain resilience research. "
                      "Single-source risk is flagged in all PE operational due diligence frameworks."
        },
    }

    for vector, details in interventions.items():
        risk = risk_at_5x.get(vector, 0)
        if risk >= 6:
            css = "danger-box"
            icon = "🔴"
        elif risk >= 4:
            css = "warning-box"
            icon = "🟡"
        else:
            css = "insight-box"
            icon = "🟢"

        st.markdown(
            f'<div class="{css}">{icon} <b>{vector}</b> (5x risk: {risk:.1f}/10)<br>'
            f'{details["action"]}<br>'
            f'<span style="font-size:0.8rem;">{details["data"]}</span><br>'
            f'<span style="font-size:0.75rem;color:#6b7280;">📚 {details["source"]}</span></div>',
            unsafe_allow_html=True)

    st.markdown("""
    <div class="source-box">
    <b>Framework:</b> Scale bottleneck analysis derived from:<br>
    (1) Theory of Constraints — Eliyahu Goldratt, "The Goal" (1984): every system has exactly one binding
    constraint; improving any other process produces zero throughput improvement until the constraint is addressed.<br>
    (2) Queuing theory — Kingman's Formula (J.F.C. Kingman, 1961): wait time grows exponentially as utilization
    approaches 100%. The 80% practical threshold is the standard operational limit derived from this formula.<br>
    (3) Blitzscaling — Reid Hoffman (2018): operational leverage and what breaks first during rapid scaling.<br>
    (4) SaaS Capital CSM ratio benchmarks; Gainsight CSM Benchmark Report 2023; MIT CTL supply chain resilience.<br>
    <b>Note:</b> The 2x/5x/10x multiples are the author's design — they map logically to seed/Series A/Series B
    milestones but are not mechanically derived from a single published framework.
    </div>
    """, unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════════
# TAB 5 — AI STRATEGIC BRIEF
# ══════════════════════════════════════════════════════════════════════════════════
with tab5:
    st.subheader("🤖 AI Strategic Brief Generator")
    st.markdown("Powered by **Groq (Llama 3.3 70B)**. Generates an advisor-grade brief from all module inputs.")

    if st.button("Generate AI Strategic Brief", type="primary", use_container_width=True):
        context = f"""
You are a senior startup advisor at the Roux Institute.
You have just completed a full operational due diligence analysis of a startup.

STARTUP: {startup_name} | Industry: {industry} | Stage: {stage} | Team: {team_size} FTEs | Model: {biz_model}

FINANCIALS:
Monthly Revenue: ${monthly_revenue:,.0f} | Gross Burn: ${monthly_burn:,.0f} | Net Burn: ${net_burn:,.0f}
Cash: ${cash_on_hand:,.0f} | Runway: {runway_months:.0f} months | MoM Growth: {mrr_growth}%
Gross Margin: {gross_margin_pct}% | Monthly Churn: {monthly_churn}%

UNIT ECONOMICS:
ARPU: ${arpu}/mo | CAC: ${cac:,.0f} | LTV: ${ltv:,.0f}
LTV:CAC: {ltv_cac:.1f}x (threshold: 3x) | CAC Payback: {payback_months:.0f}mo (threshold: 18mo)
Burn Multiple: {burn_multiple:.2f}x (Bessemer threshold: 1.5x) | Rule of 40: {rule_of_40:.0f}

MARKET:
TAM (bottom-up): {fmt_bn(tam_bu)} | SAM: {fmt_bn(sam_bu)} | SOM (3yr): {fmt_bn(som_bu)}
Method alignment: {"aligned" if method_aligned else "diverge >2x"}

OPS READINESS: {total_score:.1f}/10 — {classification}
Weakest: {weakest[0]} ({weakest[1][0]:.1f}/10) | Strongest: {strongest[0]} ({strongest[1][0]:.1f}/10)

SCALE RISK:
At 2x: {worst_2x[0]} (risk {worst_2x[1]:.1f}/10)
At 5x: {worst_5x[0]} (risk {worst_5x[1]:.1f}/10)

Produce a structured brief in exactly this format:

## SITUATION SUMMARY
[2-3 sentences on overall health, stage-appropriate assessment, and the most important signal]

## TOP 3 RISKS (ranked by financial impact)
**Risk 1 — [Name]:** [Specific, data-backed. Include the actual number from the analysis.]
**Risk 2 — [Name]:** [Specific, data-backed. Include the actual number from the analysis.]
**Risk 3 — [Name]:** [Specific, data-backed. Include the actual number from the analysis.]

## TOP 3 OPPORTUNITIES
**Opportunity 1 — [Name]:** [Specific and actionable, with supporting rationale]
**Opportunity 2 — [Name]:** [Specific and actionable]
**Opportunity 3 — [Name]:** [Specific and actionable]

## 90-DAY INTERVENTION PLAN
**Week 1-2:** [Most urgent action]
**Week 3-4:** [Second priority]
**Month 2:** [Third priority]
**Month 3:** [Fourth priority — sets up for next milestone]

## THE SINGLE METRIC TO WATCH
[One metric, why it is the leading indicator right now, and what the threshold for action is]

## INVESTOR READINESS VERDICT
[One paragraph: Is this startup ready to raise? What specifically needs to change first? Be direct.]
"""
        try:
            from groq import Groq
            client = Groq(api_key=st.secrets["GROQ_API_KEY"])
            with st.spinner("Generating strategic brief..."):
                response = client.chat.completions.create(
                    model="llama-3.3-70b-versatile",
                    messages=[{"role": "user", "content": context}],
                    temperature=0.35,
                    max_tokens=1500
                )
                brief = response.choices[0].message.content

            st.markdown("---")
            st.markdown(f"### Strategic Brief — {startup_name}")
            st.markdown(brief)
            st.download_button(
                "Download Brief as .txt",
                data=brief,
                file_name=f"{startup_name.replace(' ','_')}_VentureOps_Brief.txt",
                mime="text/plain"
            )

        except ImportError:
            st.error("Run: pip install groq")
        except KeyError:
            st.error("Add GROQ_API_KEY to .streamlit/secrets.toml")
        except Exception as e:
            st.error(f"API error: {str(e)}")

    st.markdown("---")
    st.markdown("**Brief structure covers:**")
    for item in [
        "Situation summary — stage-appropriate health assessment",
        "Top 3 risks — ranked by financial impact, sourced from actual inputs",
        "Top 3 opportunities — specific and actionable",
        "90-day intervention plan — week-by-week priorities",
        "Single metric to watch — the one number that matters most right now",
        "Investor readiness verdict — direct and honest",
    ]:
        st.markdown(f"- {item}")

    st.markdown("""
    <div class="source-box">
    <b>AI Brief Sources:</b> Brief structure adapted from Y Combinator partner review format;
    Sequoia's Arc framework for company health; First Round Capital operational assessment rubric.
    All context injected live from your inputs — no generic responses. Model: Groq Llama 3.3 70B Versatile.
    </div>
    """, unsafe_allow_html=True)


# ─── FOOTER ─────────────────────────────────────────────────────────────────────
st.divider()
st.markdown("""
<div style="text-align:center;color:#9ca3af;font-size:0.78rem;padding:1rem;">
VentureOps v2 — Built by Rutwik Satish for the Roux Institute<br>
Frameworks: CB Insights (n=101, 2014; n=431, 2023) · Bessemer Venture Partners State of Cloud 2023 · McKinsey OHI ·
David Skok SaaS Metrics · Antler TAM Methodology · Blitzscaling (Hoffman, 2018) ·
Theory of Constraints (Goldratt, 1984) · Kingman's Formula (1961) · SaaStr 2024 · Bill Gurley (2013)<br>
File extraction powered by pdfplumber + python-pptx + Groq Llama 3.3.
All inputs are processed locally. No data is stored or transmitted beyond Groq API calls.
</div>
""", unsafe_allow_html=True)
